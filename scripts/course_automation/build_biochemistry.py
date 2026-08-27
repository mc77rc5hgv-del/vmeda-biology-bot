"""Build the complete VMEDA biochemistry course from the supplied source set."""

from __future__ import annotations

import hashlib
import html
import json
import re
import shutil
import sys
from collections import Counter
from pathlib import Path

from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

MAX_CONTENT = 3300


def norm(value: str) -> str:
    value = value.replace("\u00a0", " ").replace("\r", "\n")
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def safe(value: str) -> str:
    return html.escape(norm(value), quote=False)


def chunks(text: str, limit: int = MAX_CONTENT) -> list[str]:
    paragraphs = [p.strip() for p in re.split(r"\n+", norm(text)) if p.strip()]
    result, current = [], ""
    for paragraph in paragraphs:
        pieces = [paragraph[i:i + limit] for i in range(0, len(paragraph), limit)]
        for piece in pieces:
            candidate = f"{current}\n{piece}".strip()
            if current and len(candidate) > limit:
                result.append(current)
                current = piece
            else:
                current = candidate
    if current:
        result.append(current)
    return result


def docx_text(path: Path) -> str:
    doc = Document(path)
    values = [p.text for p in doc.paragraphs if norm(p.text)]
    for table in doc.tables:
        for row in table.rows:
            value = " | ".join(norm(cell.text) for cell in row.cells)
            if norm(value):
                values.append(value)
    return norm("\n".join(values))


def docx_paragraphs(path: Path) -> list[str]:
    doc = Document(path)
    return [norm(p.text) for p in doc.paragraphs if norm(p.text)]


def pdf_pages(path: Path) -> list[str]:
    return [norm(page.extract_text() or "") for page in PdfReader(path).pages]


def lesson(lesson_id: str, title: str, body: str, source: str, media=None) -> dict:
    result = {
        "id": lesson_id,
        "title": norm(title)[:180],
        "content": f"<b>Материал курса</b>\n\n{safe(body)}",
        "sources": [source],
    }
    if media:
        result["media"] = media
    return result


def page_lessons(prefix: str, pages: list[str], source_name: str, label: str) -> list[dict]:
    result = []
    for page_no, text in enumerate(pages, 1):
        if len(text) < 40:
            continue
        page_chunks = chunks(text)
        first = re.sub(r"^\[.*?\]\s*", "", text).split("\n", 1)[0][:100]
        for part, body in enumerate(page_chunks, 1):
            suffix = f", часть {part}" if len(page_chunks) > 1 else ""
            result.append(lesson(
                f"{prefix}_p{page_no}_{part}",
                f"{label}: стр. {page_no}{suffix} — {first}", body,
                f"{source_name}, стр. {page_no}",
            ))
    return result


def numbered_lessons(prefix: str, text: str, source_name: str, label: str) -> list[dict]:
    pattern = re.compile(r"(?m)^(?:(?:ВОПРОС|БИЛЕТ)\s*№?\s*)?(\d{1,3})[\).!:–-]?\s*(?=\S)", re.I)
    matches = list(pattern.finditer(text))
    if len(matches) < 3:
        return [lesson(f"{prefix}_{i}", f"{label}, часть {i}", body, source_name)
                for i, body in enumerate(chunks(text), 1)]
    result = []
    for index, match in enumerate(matches):
        end = matches[index + 1].start() if index + 1 < len(matches) else len(text)
        block = norm(text[match.start():end])
        number = match.group(1)
        title_line = block.split("\n", 1)[0][:130]
        for part, body in enumerate(chunks(block), 1):
            extra = f", часть {part}" if len(chunks(block)) > 1 else ""
            result.append(lesson(f"{prefix}_{index + 1}_{part}", f"{label} {number}{extra}: {title_line}", body, source_name))
    return result


def ranged_question_lessons(prefix: str, text: str, source_name: str, label: str, numbers: range, marker: str = r"[)!]") -> list[dict]:
    wanted = set(numbers)
    matches = [match for match in re.finditer(rf"(?m)^(\d{{1,3}}){marker}\s*(?=\S)", text)
               if int(match.group(1)) in wanted]
    # A number can recur inside an answer; the first ordered occurrence is the question heading.
    selected, last = [], -1
    for number in numbers:
        match = next((item for item in matches if int(item.group(1)) == number and item.start() > last), None)
        if match:
            selected.append(match)
            last = match.start()
    result = []
    for index, match in enumerate(selected):
        end = selected[index + 1].start() if index + 1 < len(selected) else len(text)
        block = norm(text[match.start():end])
        parts = chunks(block)
        for part, body in enumerate(parts, 1):
            suffix = f", часть {part}" if len(parts) > 1 else ""
            result.append(lesson(f"{prefix}_{match.group(1)}_{part}", f"{label} {match.group(1)}{suffix}: {block.splitlines()[0][:130]}", body, source_name))
    return result


def multiple_choice_lessons(prefix: str, path: Path, label: str) -> list[dict]:
    values = docx_paragraphs(path)
    option = re.compile(r"^[а-яёa-z][).]\s*", re.I)
    groups, current = [], []
    for value in values:
        if option.match(value):
            if current:
                current.append(value)
            continue
        if current:
            groups.append(current)
        current = [value]
    if current:
        groups.append(current)
    result = []
    for number, group in enumerate(groups, 1):
        body = "\n".join(group)
        if len(body) < 20:
            continue
        result.append(lesson(f"{prefix}_{number}", f"{label} {number}: {group[0][:135]}", body, path.name))
    return result


def ticket_lessons(path: Path) -> list[dict]:
    text = docx_text(path)
    matches = list(re.finditer(r"(?mi)^Билет\s+(\d+)\s*$", text))
    result = []
    for index, match in enumerate(matches):
        end = matches[index + 1].start() if index + 1 < len(matches) else len(text)
        block = norm(text[match.start():end])
        for part, body in enumerate(chunks(block), 1):
            suffix = f", часть {part}" if len(chunks(block)) > 1 else ""
            result.append(lesson(f"exam_ticket_{match.group(1)}_{part}", f"Экзаменационный билет {match.group(1)}{suffix}", body, path.name))
    return result


def exam_question_lessons(path: Path) -> list[dict]:
    paragraphs = docx_paragraphs(path)
    start = max(i for i, value in enumerate(paragraphs) if value.startswith("Экзаменационные вопросы по дисциплине"))
    category = "Общие вопросы"
    result = []
    for value in paragraphs[start + 1:]:
        category_match = re.fullmatch(r"\d+\.\s*(.+)", value)
        if category_match:
            category = category_match.group(1).strip().capitalize()
            continue
        if len(value) < 10 or value.startswith(("Обсуждены", "Протокол", "Уточнено")):
            continue
        number = len(result) + 1
        result.append(lesson(f"exam_q_{number}", f"{number}. {value[:145]}", f"Раздел: {category}\n\n{value}", path.name))
    return result


def practical_exam_lessons(path: Path) -> list[dict]:
    values = docx_paragraphs(path)
    ignored = ("Военно-медицинская", "Кафедра", "УТВЕРЖДАЮ", "Заведующий", "член-корреспондент", "А. Иванов", "Обсуждены", "Протокол", "Санкт-Петербург", "Вопросы к практической")
    questions = [value for value in values if len(value) >= 12 and not value.startswith(ignored) and not re.fullmatch(r"[«_() .\dА-Яа-я]+", value)]
    return [lesson(f"exam_pr_{i}", f"Практический вопрос {i}: {value[:130]}", value, path.name)
            for i, value in enumerate(questions, 1)]


def ppt_lessons(path: Path, assets: Path) -> list[dict]:
    prs = Presentation(path)
    slide_records = []
    for slide_no, slide in enumerate(prs.slides, 1):
        texts, media = [], []
        for shape_no, shape in enumerate(slide.shapes, 1):
            if hasattr(shape, "text") and norm(shape.text):
                texts.append(norm(shape.text))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                if image.size[0] >= 500 and image.size[1] >= 300:
                    target = assets / f"lecture_s{slide_no}_{shape_no}.{image.ext}"
                    target.write_bytes(image.blob)
                    media.append({"path": f"generated_assets/biochemistry/{target.name}", "caption": f"Иллюстрация из вводной лекции, слайд {slide_no}"})
        body = norm("\n".join(texts))
        if body:
            slide_records.append((slide_no, body, media))
    result = []
    for group_no, start in enumerate(range(0, len(slide_records), 5), 1):
        group = slide_records[start:start + 5]
        first_slide, last_slide = group[0][0], group[-1][0]
        body = "\n\n".join(f"Слайд {number}. {text}" for number, text, _ in group)
        media = [item for _, _, items in group for item in items]
        for part, value in enumerate(chunks(body), 1):
            suffix = f", часть {part}" if len(chunks(body)) > 1 else ""
            result.append(lesson(f"intro_g{group_no}_{part}", f"Вводная лекция: слайды {first_slide}–{last_slide}{suffix}", value, f"Л.1 - Вводная лекция.ppt, слайды {first_slide}–{last_slide}", media if part == 1 else None))
    return result


def main(repo: Path) -> None:
    workspace = repo / ".course-automation" / "biochemistry"
    sources = workspace / "sources"
    assets = repo / "generated_assets" / "biochemistry"
    assets.mkdir(parents=True, exist_ok=True)
    for old in assets.iterdir():
        if old.is_file():
            old.unlink()

    sections = []
    sections.append({"id": "introduction", "title": "Введение в биохимию", "lessons": ppt_lessons(sources / "Л.1 - Вводная лекция.pptx", assets)})
    sections.append({"id": "core_course", "title": "Основной курс", "lessons": page_lessons("core", pdf_pages(sources / "учебное пособие.pdf"), "учебное пособие.pdf", "Основы биохимии")})
    note_lessons = page_lessons("notes", pdf_pages(sources / "ВСЯ БИОХИМИЯ.pdf"), "ВСЯ БИОХИМИЯ.pdf", "Полный курс")
    seen_note_texts = set()
    note_lessons = [item for item in note_lessons if not (item["content"] in seen_note_texts or seen_note_texts.add(item["content"]))]
    sections.append({"id": "complete_notes", "title": "Полный конспект", "lessons": note_lessons})
    sections.append({"id": "practicum", "title": "Практические и лабораторные занятия", "lessons": page_lessons("practice", pdf_pages(sources / "ПРАКТИКУМ-АЛЕКС.pdf"), "ПРАКТИКУМ-АЛЕКС.pdf", "Практикум")})

    control_sources = [
        ("control_2", "Контрольная работа №2", "БХ КР2.docx"),
        ("control_2_boundary", "Второй рубеж: вопросы 15–21", "БХ 2 рубеж 15-21.docx"),
    ]
    sections.append({"id": "control_1", "title": "Контрольная работа №1", "lessons": ranged_question_lessons("control_1", docx_text(sources / "бх кр1.docx"), "бх кр1.docx", "Вопрос", range(1, 9), marker=r"!")})
    sections.append({"id": "control_1_alt", "title": "Контрольная работа №1 — тестовый вариант", "lessons": multiple_choice_lessons("control_1_alt", sources / "Первая кр по бх.docx", "Задание")})
    for section_id, title, filename in control_sources:
        sections.append({"id": section_id, "title": title, "lessons": numbered_lessons(section_id, docx_text(sources / filename), filename, "Вопрос")})
    sections.append({"id": "control_3", "title": "Контрольная работа №3: вопросы 23–30", "lessons": ranged_question_lessons("control_3", docx_text(sources / "БХ(3)23-30.docx"), "БХ(3)23-30.docx", "Вопрос", range(23, 31))})
    sections.append({"id": "tests", "title": "Все тесты по биохимии", "lessons": multiple_choice_lessons("test", sources / "тут все тесты по бх.docx", "Тест")})
    sections.append({"id": "credit", "title": "Зачёт", "lessons": page_lessons("credit", pdf_pages(sources / "c_биохимия зачет все вопросы.pdf"), "c_биохимия зачет все вопросы.pdf", "Зачёт")})
    sections.append({"id": "exam_tickets", "title": "Экзамен — билеты с ответами", "lessons": ticket_lessons(sources / "BKh_EKZ_BILETY.docx")})
    sections.append({"id": "exam_questions", "title": "Экзамен — перечень теоретических вопросов", "lessons": exam_question_lessons(sources / "экзаменационные_вопросы_БХ_ЛД_7c9f668527fe7e2b57547f1efd175d19.docx")})
    sections.append({"id": "exam_practical", "title": "Экзамен — практическая часть", "lessons": practical_exam_lessons(sources / "Вопросы_Практическая_часть_экзамена_Ст.docx")})

    # The large textbook is a reference source. Its extracted contents are exposed as a navigable map.
    textbook_pages = pdf_pages(sources / "Северин_учебник_9bea3e614f2611bd4bfa3da857615d0b.pdf")
    map_text = "\n".join(textbook_pages[:8])
    sections.append({"id": "textbook", "title": "Учебник Северина — содержание и навигация", "lessons": [
        lesson(f"textbook_map_{i}", f"Учебник Северина: содержание, часть {i}", body, "Северин_учебник_9bea3e614f2611bd4bfa3da857615d0b.pdf, стр. 1–8")
        for i, body in enumerate(chunks(map_text), 1)
    ]})

    course = {
        "id": "biochemistry",
        "course": 2,
        "title": "Биохимия",
        "emoji": "🧬",
        "description": "Полный курс биохимии по материалам ВМедА: теория, практикум, контрольные работы, тесты, зачёт и экзамен.",
        "sections": sections,
    }
    output = repo / "generated_courses" / "biochemistry.json"
    output.write_text(json.dumps(course, ensure_ascii=False, indent=2), encoding="utf-8")

    originals = [p for p in sources.iterdir() if p.is_file() and p.suffix.lower() in {".ppt", ".docx", ".pdf"}]
    used = Counter(source.split(",")[0] for section in sections for item in section["lessons"] for source in item.get("sources", []))
    report = {
        "subject": "Биохимия",
        "original_source_count": len(originals),
        "section_count": len(sections),
        "lesson_count": sum(len(s["lessons"]) for s in sections),
        "media_count": sum(len(l.get("media", [])) for s in sections for l in s["lessons"]),
        "sources": [{
            "file": p.name,
            "bytes": p.stat().st_size,
            "sha256": hashlib.sha256(p.read_bytes()).hexdigest(),
            "included_lessons": used[p.name],
            "role": "reference" if p.name.startswith("Северин_") else "course",
        } for p in sorted(originals, key=lambda p: p.name.casefold())],
    }
    (workspace / "coverage_report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({k: report[k] for k in ("original_source_count", "section_count", "lesson_count", "media_count")}, ensure_ascii=False))


if __name__ == "__main__":
    main(Path(sys.argv[1]).resolve())
