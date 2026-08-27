import json
from pathlib import Path


def _pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n\n".join(f"[Страница {i + 1}]\n{page.extract_text() or ''}" for i, page in enumerate(reader.pages))


def _docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    chunks = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    for table in doc.tables:
        chunks.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
    return "\n".join(chunks)


def _pptx(path: Path) -> str:
    from pptx import Presentation
    deck = Presentation(path)
    chunks = []
    for index, slide in enumerate(deck.slides, 1):
        text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        chunks.append(f"[Слайд {index}]\n" + "\n".join(text))
    return "\n\n".join(chunks)


def extract_sources(workspace: Path) -> dict:
    source_dir = workspace / "sources"
    output_dir = workspace / "extracted"
    output_dir.mkdir(parents=True, exist_ok=True)
    records = []
    handlers = {".pdf": _pdf, ".docx": _docx, ".pptx": _pptx}
    for source in sorted(source_dir.glob("*")):
        if not source.is_file():
            continue
        try:
            if source.suffix.lower() in handlers:
                text = handlers[source.suffix.lower()](source)
            elif source.suffix.lower() in {".txt", ".md", ".csv"}:
                text = source.read_text(encoding="utf-8", errors="replace")
            else:
                continue
            target = output_dir / f"{source.name}.txt"
            target.write_text(text, encoding="utf-8")
            records.append({"source": source.name, "text": str(target), "characters": len(text)})
        except Exception as exc:
            records.append({"source": source.name, "error": str(exc)})
    report = {"files": records, "successful": sum("text" in item for item in records)}
    (workspace / "extraction_report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    return report

